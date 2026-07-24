"""Note loading, embedding, and caching pipeline.

This module is responsible for three things:

1. File ingestion — reads .txt, .md, .pdf, .docx, and .pptx files from
   the user's chosen notes folder and stores their plain-text content in
   the Notes dict.
2. Semantic embedding — encodes each note with a sentence-transformer model
   so similarity search can be performed without a round-trip to the LLM.
3. Disk cache — persists embeddings between sessions using pickle so the
   model doesn't re-encode unchanged files on every launch. The cache key
   is an MD5 hash of the file bytes so stale entries are detected automatically.

Module-level globals are used intentionally: loader is treated as a
shared-state namespace by retrieval.py and UI.py (they import loader and
read loader.Notes, loader.model, etc. directly).
"""

# loaders.py

import hashlib
import pickle
from pathlib import Path
import io
import fitz
from docx import Document
from pptx import Presentation
from routes.config import EMBEDDINGS_CACHE_FILE

# Shared state — imported by retrieval.py and UI.py
Notes = {}              # Filename → plain-text content for every loaded note
Embedding_cache = {}    # Filename → sentence-transformer tensor (in-memory)
Notes_Cache = {}        # Reserved for future use (currently unused)
model = None            # SentenceTransformer instance, set in bg_model_loading()
model_loaded = False    # True once bg_model_loading() finishes


# def has_supported_files(folder_path):
#     """Check whether a folder contains at least one supported note file.

#     Args:
#         folder_path: Absolute path to the folder to check.

#     Returns:
#         True if at least one file with a supported extension exists,
#         False otherwise.
#     """
#     for filename in os.listdir(folder_path):
#         if Path(os.path.join(folder_path, filename)).suffix.lower() in ALLOWED_EXTENSIONS:
#             return True
#     return False

# def reset_notes():
#     Notes.clear()
#     Embedding_cache.clear()
#     Notes_Cache.clear()
    
def store_notes(filename, filecontent):
    """Scan folder_path and load all supported note files into Notes.

    For each file the text content is extracted via the appropriate loader,
    the MD5 hash is compared against the on-disk embedding cache, and only
    changed files are re-encoded. Empty files are silently skipped; unreadable
    files are logged and skipped without raising.

    Args:
        folder_path: Absolute path to the folder containing the note files.
    """
    # Clear state so a folder switch doesn't mix old and new notes
    disk_cache = {}

    # Load the on-disk embedding cache so we can skip re-encoding unchanged files
    if EMBEDDINGS_CACHE_FILE.exists():
        with open(EMBEDDINGS_CACHE_FILE, "rb") as f:
            disk_cache = pickle.load(f)
    notes_changed = False  # Track whether any file is new or modified



        # Skip files with unsupported extensions
    extensions = Path(filename).suffix.lower()
    if extensions == '.pdf':
        content = pdf_loader(filecontent)
    elif extensions == '.docx':
        content = docx_loader(filecontent)
    elif extensions == '.pptx':
        content = presentation_loader(filecontent)
    elif extensions in ['.txt','.md']:
        content = filecontent.decode('utf-8', errors="ignore")
    else:
        content = "" 
    try:
        if not content:
            return
        file_hash = get_hash(filecontent)
        Notes[filename] = content   

        if filename in disk_cache and disk_cache[filename][0] == file_hash:
            # File unchanged — reuse the cached embedding tensor
            Embedding_cache[filename] = disk_cache[filename][1]
        else:
            # File is new or modified — encode and update cache
            Embedding_cache[filename] = model.encode(
                content, convert_to_tensor=True
            )
            disk_cache[filename] = (file_hash, Embedding_cache[filename])
            notes_changed = True

    except Exception as e:
        print(f"Skipping unreadable asset {filename}: {e}")

    
    # Only write to disk if something changed to avoid unnecessary I/O
    if notes_changed:
        with open(EMBEDDINGS_CACHE_FILE, "wb") as f:
            pickle.dump(disk_cache, f)


def bg_model_loading():
    """Load the embedding model and notes in a background thread.

    Intended to run in a threading.Thread so the splash screen remains
    responsive during the potentially slow model download and note-indexing
    steps. All widget updates use ui_instance.after() which is thread-safe.

    Args:
        ui_instance: The ACEUI instance whose status label should be updated.

    Returns:
        A (model, model_loaded) tuple; callers typically ignore the return
        value since both are accessible as module-level globals.
    """
    global model, model_loaded



    # Update splash screen — must use .after() because this runs off the main thread
    # ui_instance.after(
    #     0, lambda: ui_instance.update_loading_status("LOADING EMBEDDINGS ENGINE...")
    # )

    # Deferred import keeps cold-start time low when the model is already cached
    from sentence_transformers import SentenceTransformer

    model = SentenceTransformer("all-MiniLM-L6-v2")

    # Load the folder path and index notes (re-encodes only changed files)
    model_loaded = True

    # Signal the main thread to swap splash screen for main UI
    return model, model_loaded


def get_hash(filecontent):
    if  not isinstance(filecontent, bytes):
        return None
    return hashlib.md5(filecontent).hexdigest()
    """Compute the MD5 hex digest of a file's raw bytes.

    Used to detect whether a note file has changed since it was last encoded
    so the embedding cache can be invalidated selectively.

    Args:
        filepath: Absolute path to the file to hash.

    Returns:
        A lowercase hex string of the MD5 digest (32 characters).
    """
    


def pdf_loader(filecontent):
    """Extract all plain text from a PDF file using pdfplumber.

    Iterates over every page and concatenates the extracted text. Pages that
    yield None (e.g. image-only pages) are silently skipped via the ``or ""``
    fallback.

    Args:
        file_path: Absolute path to the .pdf file.

    Returns:
        Concatenated text of all pages, stripped of leading/trailing whitespace.
    """
    doc = fitz.open(stream=filecontent, filetype='pdf')
    text = ""
    for page in doc:
        text += page.get_text()
    
    return text.strip()


def docx_loader(filecontent):
    """Extract plain text from a .docx Word document.

    Reads paragraph objects from python-docx and concatenates their text.
    Formatting, images, and tables are not extracted.

    Args:
        filepath: Absolute path to the .docx file.

    Returns:
        Concatenated paragraph text, stripped of leading/trailing whitespace.
    """
    text = ""
    doc = Document(io.BytesIO(filecontent))
    for paragraph in doc.paragraphs:
        text += paragraph.text or ""
    return text.strip()


def presentation_loader(filecontent):
    """Extract plain text from a .pptx PowerPoint presentation.

    Iterates over all slides, shapes, and text-frame paragraphs. Only shapes
    that expose a text_frame attribute are read; images and charts are skipped.

    Args:
        filecontent: Absolute path to the .pptx file.

    Returns:
        Concatenated slide text, stripped of leading/trailing whitespace.
    """
    text = ""
    prs = Presentation(io.BytesIO(filecontent))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text or ""
    return text.strip()
